"""Parse .md / .pdf / .docx / .pptx files into plain text for ingest."""

from __future__ import annotations

import re
from pathlib import Path

# ------------------------------------------------------------------
# Supported formats
# ------------------------------------------------------------------
SUPPORTED = {".md", ".pdf", ".docx", ".pptx"}


# ------------------------------------------------------------------
# Per-format parsers
# ------------------------------------------------------------------
def parse_markdown(raw: bytes) -> str:
    text = raw.decode("utf-8", errors="replace")
    # Strip fenced code blocks
    text = re.sub(r"```[\s\S]*?```", "", text)
    # Strip inline code
    text = re.sub(r"`[^`]+`", "", text)
    # Strip markdown links, keep text
    text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text)
    # Strip images
    text = re.sub(r"!\[([^\]]*)\]\([^\)]+\)", "", text)
    # Strip ATX headings markers
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    # Strip blockquotes
    text = re.sub(r"^>\s+", "", text, flags=re.MULTILINE)
    # Strip horizontal rules
    text = re.sub(r"^[-*_]{3,}\s*$", "", text, flags=re.MULTILINE)
    # Collapse blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def parse_pdf(raw: bytes) -> str:
    try:
        import pdfplumber
    except ImportError:
        return "[无法解析 PDF 文件，请将文本内容复制粘贴导入]"
    import io as _io
    with pdfplumber.open(_io.BytesIO(raw)) as pdf:
        pages = []
        for page in pdf.pages:
            text = (page.extract_text() or "").strip()
            if text:
                # Collapse excessive whitespace
                text = re.sub(r"\s+", " ", text)
                pages.append(text)
    result = "\n\n".join(pages)
    if len(result) < 100:
        return "[PDF 文件内容无法自动提取（可能是扫描件），请将文本复制粘贴导入]"
    return result


def parse_docx(raw: bytes) -> str:
    try:
        import docx
    except ImportError:
        return "[无法解析 Word 文件，请将文本内容复制粘贴导入]"
    import io as _io
    doc = docx.Document(_io.BytesIO(raw))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def parse_pptx(raw: bytes) -> str:
    try:
        import pptx
    except ImportError:
        return "[无法解析 PowerPoint 文件，请将文本内容复制粘贴导入]"
    import io as _io
    prs = pptx.Presentation(_io.BytesIO(raw))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        if lines:
            slides.append(f"--- 第 {i} 页 ---\n" + "\n".join(lines))
    return "\n\n".join(slides)


# ------------------------------------------------------------------
# Main entry point (sync)
# ------------------------------------------------------------------
def parse_file(file_bytes: bytes, filename: str) -> str:
    """Parse a file by extension and return plain text."""
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED:
        raise ValueError(f"不支持的格式: {ext}，支持的格式: {', '.join(sorted(SUPPORTED))}")

    if ext == ".md":
        return parse_markdown(file_bytes)
    if ext == ".pdf":
        return parse_pdf(file_bytes)
    if ext == ".docx":
        return parse_docx(file_bytes)
    if ext == ".pptx":
        return parse_pptx(file_bytes)

    raise ValueError(f"未实现的格式: {ext}")
